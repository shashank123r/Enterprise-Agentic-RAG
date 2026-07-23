"""PPTX extractor — sync calls offloaded via run_in_executor."""

from app.core.logging import get_logger
from app.ingestion.executor import run_in_executor
from app.ingestion.extractors import (
    BaseExtractor, ExtractionResult, ImageResult, PageResult, TableResult, extractor_registry,
)

logger = get_logger(__name__)


@extractor_registry.register("application/vnd.openxmlformats-officedocument.presentationml.presentation")
class PptxExtractor(BaseExtractor):
    """Extract text, tables, and images from PPTX presentations."""

    async def extract(self) -> ExtractionResult:
        def _do() -> ExtractionResult:
            from pptx import Presentation
            prs = Presentation(str(self.file_path))
            r = ExtractionResult()
            props = prs.core_properties
            r.metadata.update({"title": props.title or "", "author": props.author or "", "created": str(props.created) if props.created else "", "modified": str(props.modified) if props.modified else "", "slide_count": len(prs.slides)})
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text: list[str] = []
                slide_title: str | None = None
                table_idx = 0
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if not text: continue
                        if slide_title is None: slide_title = text
                        slide_text.append(text)
                    if shape.has_table:
                        tbl = shape.table
                        headers = [cell.text.strip() for cell in tbl.rows[0].cells]
                        rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows[1:]]
                        r.tables.append(TableResult(page_number=slide_num, table_index=table_idx, headers=headers, rows=rows, caption=f"Slide {slide_num} table {table_idx + 1}"))
                        table_idx += 1
                if slide_text:
                    r.pages.append(PageResult(page_number=slide_num, text="\n".join(slide_text), section_title=slide_title))
                image_count = 0
                for shape in slide.shapes:
                    if shape.shape_type == 13:
                        image_count += 1
                        img = shape.image
                        r.images.append(ImageResult(page_number=slide_num, image_index=image_count, image_data=img.blob, format=img.content_type.split("/")[-1] or "png"))
            r.text = "\n\n".join(p.text for p in r.pages if p.text.strip())
            return r
        return await run_in_executor(_do)

    async def extract_text(self) -> str:
        def _do() -> str:
            from pptx import Presentation
            prs = Presentation(str(self.file_path))
            texts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text.strip())
            return "\n".join(t for t in texts if t)
        return await run_in_executor(_do)

    async def extract_metadata(self) -> dict:
        def _do() -> dict:
            from pptx import Presentation
            prs = Presentation(str(self.file_path))
            props = prs.core_properties
            return {"title": props.title or "", "author": props.author or "", "created": str(props.created) if props.created else "", "modified": str(props.modified) if props.modified else "", "slide_count": len(prs.slides)}
        return await run_in_executor(_do)
